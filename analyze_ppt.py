
from pptx import Presentation
from pathlib import Path

ppt_path = Path("PPT/Fortune3_HandloomHackathon2026.pptx")
prs = Presentation(ppt_path)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- Slide {i} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(f"  - {shape.text[:200]}")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    print(f"    > {p.text[:200]}")
