
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ppt_path = Path("PPT/Fortune3_HandloomHackathon2026.pptx")
prs = Presentation(ppt_path)

slide_width = prs.slide_width
slide_height = prs.slide_height

# ─────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────
def clear_text_frame(tf):
    for p in list(tf.paragraphs):
        for r in list(p.runs):
            r.text = ""
    if len(tf.paragraphs) > 0:
        return tf.paragraphs[0]
    return tf.add_paragraph()

def set_text(shape, text, font_size=24, bold=False, color=RGBColor(0x20, 0x1A, 0x17), align=PP_ALIGN.LEFT):
    if shape.has_text_frame:
        p = clear_text_frame(shape.text_frame)
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return p
    return None

def add_paragraph(shape, text, font_size=20, bold=False, color=RGBColor(0x3D, 0x34, 0x2F), space_before=Pt(8)):
    if shape.has_text_frame:
        p = shape.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = space_before
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return p
    return None

# Colors
PRIMARY = RGBColor(0x5C, 0x40, 0x24)     # warm brown
ACCENT = RGBColor(0xA6, 0x7B, 0x4D)      # textile tan
INK = RGBColor(0x20, 0x1A, 0x17)         # dark ink
MUTED = RGBColor(0x5A, 0x4F, 0x4A)       # muted text
BG_CREAM = RGBColor(0xFA, 0xF5, 0xEE)    # cream background

# ─────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        if "2026" in shape.text:
            # Replace with proper title
            p = clear_text_frame(shape.text_frame)
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "SAATHI"
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            p2 = shape.text_frame.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(16)
            r2 = p2.add_run()
            r2.text = "AI Companion for Handloom Weavers"
            r2.font.size = Pt(28)
            r2.font.color.rgb = MUTED

            p3 = shape.text_frame.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(32)
            r3 = p3.add_run()
            r3.text = "Income Stability & Demand Forecasting Tools"
            r3.font.size = Pt(20)
            r3.font.color.rgb = ACCENT

            p4 = shape.text_frame.add_paragraph()
            p4.alignment = PP_ALIGN.CENTER
            p4.space_before = Pt(80)
            r4 = p4.add_run()
            r4.text = "Problem Statement 4.2 | 2026"
            r4.font.size = Pt(16)
            r4.font.color.rgb = MUTED

# ─────────────────────────────────────────────────────────────
# SLIDE 2 — Team (THE MAIN FIX)
# ─────────────────────────────────────────────────────────────
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        # Clear and rebuild
        tf = shape.text_frame
        tf.word_wrap = True
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                r.text = ""

        # Heading: TEAM FORTUNE III
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = "TEAM FORTUNE III"
        r0.font.size = Pt(40)
        r0.font.bold = True
        r0.font.color.rgb = PRIMARY

        # Separator
        p_sep = tf.add_paragraph()
        p_sep.alignment = PP_ALIGN.CENTER
        p_sep.space_before = Pt(16)
        r_sep = p_sep.add_run()
        r_sep.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
        r_sep.font.size = Pt(16)
        r_sep.font.color.rgb = ACCENT

        # Kathan
        p1 = tf.add_paragraph()
        p1.space_before = Pt(24)
        p1.alignment = PP_ALIGN.CENTER
        r1a = p1.add_run()
        r1a.text = "👤  Kathan Patel"
        r1a.font.size = Pt(24)
        r1a.font.bold = True
        r1a.font.color.rgb = INK

        p1b = tf.add_paragraph()
        p1b.alignment = PP_ALIGN.CENTER
        p1b.space_before = Pt(4)
        r1b = p1b.add_run()
        r1b.text = "Project Lead • AI/ML • Full Stack Development"
        r1b.font.size = Pt(16)
        r1b.font.color.rgb = MUTED

        # Varun
        p2 = tf.add_paragraph()
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.CENTER
        r2a = p2.add_run()
        r2a.text = "👤  Varun Kushvaha"
        r2a.font.size = Pt(24)
        r2a.font.bold = True
        r2a.font.color.rgb = INK

        p2b = tf.add_paragraph()
        p2b.alignment = PP_ALIGN.CENTER
        p2b.space_before = Pt(4)
        r2b = p2b.add_run()
        r2b.text = "Frontend Development • UI/UX"
        r2b.font.size = Pt(16)
        r2b.font.color.rgb = MUTED

        # Dhruv
        p3 = tf.add_paragraph()
        p3.space_before = Pt(20)
        p3.alignment = PP_ALIGN.CENTER
        r3a = p3.add_run()
        r3a.text = "👤  Dhruv Gohel"
        r3a.font.size = Pt(24)
        r3a.font.bold = True
        r3a.font.color.rgb = INK

        p3b = tf.add_paragraph()
        p3b.alignment = PP_ALIGN.CENTER
        p3b.space_before = Pt(4)
        r3b = p3b.add_run()
        r3b.text = "Backend Development • Data Engineering"
        r3b.font.size = Pt(16)
        r3b.font.color.rgb = MUTED

        # Separator
        p_sep2 = tf.add_paragraph()
        p_sep2.alignment = PP_ALIGN.CENTER
        p_sep2.space_before = Pt(28)
        r_sep2 = p_sep2.add_run()
        r_sep2.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
        r_sep2.font.size = Pt(16)
        r_sep2.font.color.rgb = ACCENT

        # University info
        p_uni1 = tf.add_paragraph()
        p_uni1.alignment = PP_ALIGN.CENTER
        p_uni1.space_before = Pt(24)
        r_uni1 = p_uni1.add_run()
        r_uni1.text = "🏛  Gujarat Technological University"
        r_uni1.font.size = Pt(20)
        r_uni1.font.bold = True
        r_uni1.font.color.rgb = INK

        p_uni2 = tf.add_paragraph()
        p_uni2.alignment = PP_ALIGN.CENTER
        p_uni2.space_before = Pt(4)
        r_uni2 = p_uni2.add_run()
        r_uni2.text = "School of Engineering & Technology"
        r_uni2.font.size = Pt(16)
        r_uni2.font.color.rgb = MUTED

        p_ps = tf.add_paragraph()
        p_ps.alignment = PP_ALIGN.CENTER
        p_ps.space_before = Pt(20)
        r_ps = p_ps.add_run()
        r_ps.text = "Problem Statement 4.2 • Income Stability & Demand Forecasting Tools"
        r_ps.font.size = Pt(14)
        r_ps.font.color.rgb = ACCENT
        r_ps.font.italic = True

# ─────────────────────────────────────────────────────────────
# SLIDE 3 — Challenges / Problem
# ─────────────────────────────────────────────────────────────
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        txt = shape.text
        if "Challenges Faced" in txt or "The Problem" in txt:
            # Heading
            # Keep as is, but maybe clean up
            pass
        elif "Seasonal Demand" in txt and "Delayed Payments" in txt:
            # This is the main problem content - make it cleaner
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""

            # Subtitle: The Problem (clearer)
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = "🧵  Why Weavers Struggle"
            r0.font.size = Pt(28)
            r0.font.bold = True
            r0.font.color.rgb = PRIMARY

            # Items
            items = [
                ("📈  Unpredictable Seasonal Demand", "Demand shifts with festivals and trends — no way to plan ahead."),
                ("💰  Delayed Payments", "Payments arrive weeks after delivery — cashflow dries up."),
                ("📦  Unsold Inventory Risk", "Overproduction from guesswork leaves unsold stock tying up capital."),
                ("📊  No Data-Driven Decisions", "Production choices rely on personal experience, not market signals."),
                ("📉  Chronic Income Instability", "Uncertain demand + cashflow gaps make long-term planning impossible."),
            ]
            for i, (title, desc) in enumerate(items):
                pi = tf.add_paragraph()
                pi.space_before = Pt(14 if i == 0 else 12)
                ri = pi.add_run()
                ri.text = title
                ri.font.size = Pt(20)
                ri.font.bold = True
                ri.font.color.rgb = INK

                pid = tf.add_paragraph()
                pid.space_before = Pt(2)
                rid = pid.add_run()
                rid.text = desc
                rid.font.size = Pt(14)
                rid.font.color.rgb = MUTED

            # Gap statement
            pgap = tf.add_paragraph()
            pgap.space_before = Pt(24)
            rgap = pgap.add_run()
            rgap.text = "⚠️  Current tools forecast demand — they don't tell weavers what to weave next week."
            rgap.font.size = Pt(16)
            rgap.font.bold = True
            rgap.font.color.rgb = ACCENT

# ─────────────────────────────────────────────────────────────
# SLIDE 4 — Solution
# ─────────────────────────────────────────────────────────────
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame:
        txt = shape.text
        if "SAATHI transforms" in txt:
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""

            # Core idea
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = "✨  What SAATHI Does"
            r0.font.size = Pt(28)
            r0.font.bold = True
            r0.font.color.rgb = PRIMARY

            p1 = tf.add_paragraph()
            p1.space_before = Pt(12)
            r1 = p1.add_run()
            r1.text = (
                "Turns 5 years of demand history into a trusted weekly plan. "
                "Instead of showing graphs, SAATHI answers one question:"
            )
            r1.font.size = Pt(16)
            r1.font.color.rgb = MUTED

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)
            r2 = p2.add_run()
            r2.text = "“What should I do this week to earn more money?”"
            r2.font.size = Pt(22)
            r2.font.bold = True
            r2.font.italic = True
            r2.font.color.rgb = PRIMARY

            # Feature grid text
            p3 = tf.add_paragraph()
            p3.space_before = Pt(28)
            r3 = p3.add_run()
            r3.text = "Capabilities that build trust:"
            r3.font.size = Pt(18)
            r3.font.bold = True
            r3.font.color.rgb = INK

        elif "Demand Forecasting" in txt and "Cashflow" in txt:
            # Clean up the feature list
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""

            features = [
                "📈  Demand Forecasting  →  Confidence-weighted weekly predictions",
                "💰  Cashflow Estimation  →  Loan guidance & payment windows",
                "🧠  Explainable AI  →  Clear “why” for every recommendation",
                "🌐  Multilingual  →  English • Hindi • Gujarati (no mixed UI)",
                "🎤  Voice Assistant  →  Speech-first for low-literacy weavers",
                "📋  Weekly Plan  →  Day-by-day weaving tasks + buy/sell dates",
            ]
            for i, f in enumerate(features):
                pi = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if i > 0:
                    pi.space_before = Pt(10)
                ri = pi.add_run()
                ri.text = f
                ri.font.size = Pt(16)
                ri.font.color.rgb = INK

# ─────────────────────────────────────────────────────────────
# SLIDE 5 — Prototype (add a text description)
# ─────────────────────────────────────────────────────────────
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        txt = shape.text
        if "Prototype / Product" in txt:
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run()
            r0.text = "📱  Product Snapshot"
            r0.font.size = Pt(32)
            r0.font.bold = True
            r0.font.color.rgb = PRIMARY

            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(20)
            r1 = p1.add_run()
            r1.text = "Mobile-first • Warm textile aesthetic • Zero-jargon UI"
            r1.font.size = Pt(18)
            r1.font.color.rgb = MUTED

            # Add a description block
            p2 = tf.add_paragraph()
            p2.space_before = Pt(40)
            r2 = p2.add_run()
            r2.text = "What's on the home screen:"
            r2.font.size = Pt(18)
            r2.font.bold = True
            r2.font.color.rgb = INK

            bullets = [
                "✦  Hero card: “Weave 6 Cotton Sarees · High Demand · High Confidence”",
                "✦  Expected weekly income + cash health + savings estimate",
                "✦  Buy material now · Sell window 13 Jul – 20 Jul",
                "✦  My Progress: income trend, reliability, orders, financial health",
                "✦  One-tap AI assistant + voice input for instant questions",
            ]
            for i, b in enumerate(bullets):
                pb = tf.add_paragraph()
                pb.space_before = Pt(8)
                rb = pb.add_run()
                rb.text = b
                rb.font.size = Pt(15)
                rb.font.color.rgb = MUTED

# ─────────────────────────────────────────────────────────────
# SLIDE 6 — Benefits & Future
# ─────────────────────────────────────────────────────────────
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        txt = shape.text
        if "Experience SAATHI" in txt or "Live Prototype" in txt:
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""

            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run()
            r0.text = "🚀  Impact & What's Next"
            r0.font.size = Pt(32)
            r0.font.bold = True
            r0.font.color.rgb = PRIMARY

            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(16)
            r1 = p1.add_run()
            r1.text = "Immediate Benefits for Weavers"
            r1.font.size = Pt(20)
            r1.font.bold = True
            r1.font.color.rgb = INK

            benefits = [
                "✅  Confidence to plan production without guesswork",
                "✅  Lower inventory risk — produce only what the market wants",
                "✅  Steadier income across seasons and festivals",
                "✅  Simple actions in the weaver's own language",
                "✅  Trust from explained reasoning, not black-box numbers",
            ]
            for b in benefits:
                pb = tf.add_paragraph()
                pb.space_before = Pt(8)
                rb = pb.add_run()
                rb.text = b
                rb.font.size = Pt(16)
                rb.font.color.rgb = MUTED

            p_fut = tf.add_paragraph()
            p_fut.space_before = Pt(24)
            p_fut.alignment = PP_ALIGN.CENTER
            r_fut = p_fut.add_run()
            r_fut.text = "🔮  Roadmap"
            r_fut.font.size = Pt(20)
            r_fut.font.bold = True
            r_fut.font.color.rgb = ACCENT

            roadmap = [
                "Phase 1 · AI Companion (current) → Weekly plan, cashflow, voice chat",
                "Phase 2 · Market Linkage → Buyer matching, order tracking, payments",
                "Phase 3 · Cluster Network → Shared procurement, group pricing, skilling",
            ]
            for rm in roadmap:
                prm = tf.add_paragraph()
                prm.space_before = Pt(6)
                prm.alignment = PP_ALIGN.CENTER
                rrm = prm.add_run()
                rrm.text = rm
                rrm.font.size = Pt(14)
                rrm.font.color.rgb = MUTED

            # Live prototype
            plive = tf.add_paragraph()
            plive.alignment = PP_ALIGN.CENTER
            plive.space_before = Pt(28)
            rlive = plive.add_run()
            rlive.text = "🔗  Try the Live Prototype"
            rlive.font.size = Pt(22)
            rlive.font.bold = True
            rlive.font.color.rgb = PRIMARY

            purl = tf.add_paragraph()
            purl.alignment = PP_ALIGN.CENTER
            purl.space_before = Pt(6)
            rurl = purl.add_run()
            rurl.text = "handloom-ai-app.onrender.com"
            rurl.font.size = Pt(20)
            rurl.font.underline = True
            rurl.font.color.rgb = ACCENT

# ─────────────────────────────────────────────────────────────
# SLIDE 7 — Thank you
# ─────────────────────────────────────────────────────────────
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        txt = shape.text
        if "Thank you" in txt:
            tf = shape.text_frame
            tf.word_wrap = True
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    r.text = ""

            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run()
            r0.text = "धन्यवाद  🙏  Thank You"
            r0.font.size = Pt(54)
            r0.font.bold = True
            r0.font.color.rgb = PRIMARY

            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(32)
            r1 = p1.add_run()
            r1.text = "Team Fortune III"
            r1.font.size = Pt(28)
            r1.font.bold = True
            r1.font.color.rgb = INK

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(12)
            r2 = p2.add_run()
            r2.text = "Kathan Patel  •  Varun Kushvaha  •  Dhruv Gohel"
            r2.font.size = Pt(18)
            r2.font.color.rgb = MUTED

            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(12)
            r3 = p3.add_run()
            r3.text = "GTU · School of Engineering & Technology"
            r3.font.size = Pt(16)
            r3.font.color.rgb = MUTED

            p4 = tf.add_paragraph()
            p4.alignment = PP_ALIGN.CENTER
            p4.space_before = Pt(48)
            r4 = p4.add_run()
            r4.text = "Transforming Forecasts into Better Decisions"
            r4.font.size = Pt(20)
            r4.font.italic = True
            r4.font.color.rgb = ACCENT

# ─────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────
out_path = Path("PPT/Fortune3_HandloomHackathon2026_IMPROVED.pptx")
prs.save(out_path)
print(f"✅  Saved improved PPT to {out_path}")
print(f"   Total slides: {len(prs.slides)}")
